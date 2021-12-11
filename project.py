from io import FileIO
import io
import os
from wand.image import Image
from pptx.util import Inches 
from pptx import Presentation   
from tqdm import tqdm
import gzip

  
def addWatermark()->FileIO:
    print("Adding watermark to the images just hold on....")
    try:
        for file in tqdm(os.listdir()):
            if file.endswith('.jpg') or file.endswith('.jpeg'):
            # Import the image
                with Image(filename =file) as image:
                    
                    # Extract the image's width and height
                    width, height = image.size
                    
                    with Image(filename ='nike_black.png') as logo:
                        
                        # Resize the logo image according to the image size
                        logo.resize(width//3, height//8)
                        
                        # Clone the image in order to process
                        with image.clone() as watermark:
                            # Resize the watermark
                            watermark.watermark(logo)
                            # Save the image
                            watermark.save(filename='watermarked_'+file)
                        
        print('Congratulations!!! Nike logo has been successfully added into every images\n')
    except Exception as e:
        print('Oops :-( There seems to be error while adding watermark to the images')
        
    
        
                        
def createSlides()->FileIO:
    # Creating presentation object
    root = Presentation()
    print("Creating slides just hold on ....")
    n=1
    try:
        for file in tqdm(os.listdir()):
            if file.startswith('watermarked_'):
                # Creating slide layout
                first_slide_layout = root.slide_layouts[1] 
                slide = root.slides.add_slide(first_slide_layout)
                shapes = slide.shapes
                
                #Adding title or heading to the slide
                title_shape = shapes.title
                title_shape.text = "Slides created using python-pptx"
                
                #Adding sub-title with border to the slide
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = f"This is subtitle of image {n}"
                n += 1
                    
                #Adding image to the slide
                with Image(filename = file) as watermarked_image:
                    
                    #Maintianing the aspect ratio of the image
                    width, height = watermarked_image.size
                    ratio = height/width
                    new_width = width / 25
                    new_height = int(new_width * ratio)
                    watermarked_image.resize(int(new_width), new_height)
                    
                    #Without byteIO the image is not saved or produced an eror
                    buf = io.BytesIO()
                    watermarked_image.save(buf)
                    buf.seek(0)
                    
                    # Add the watermarked image to the slide
                    slide.shapes.add_picture(buf, Inches(1), Inches(2.5))
                    
                    #saving the presentation
                    root.save("Output.pptx")
        print("Congratulations !!! slides has been created successfully")
    except Exception as e:
        print('Oops :-( There seems to be error while creating slides')

#calling the functions to create watermark       
addWatermark()
#calling the function to create slides
createSlides()